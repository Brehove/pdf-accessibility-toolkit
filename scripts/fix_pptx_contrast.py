#!/usr/bin/env python3
"""
Fix PPTX Contrast Issues for Accessibility
Analyzes and fixes text contrast issues to meet WCAG 2.1 AA requirements.

Usage:
    python fix_pptx_contrast.py input.pptx
    python fix_pptx_contrast.py input.pptx --analyze  # Just show issues, don't fix
    python fix_pptx_contrast.py *.pptx  # Batch process
"""

import argparse
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
import re


# WCAG AA Contrast Requirements
CONTRAST_NORMAL = 4.5  # Normal text
CONTRAST_LARGE = 3.0   # Large text (18pt+ or 14pt+ bold)

# Colors for fixes
DARK_BLUE = (0x2E, 0x75, 0xB6)  # #2E75B6 - good contrast with white
DARK_GRAY = (0x40, 0x40, 0x40)  # #404040 - good contrast with white backgrounds
WHITE = (0xFF, 0xFF, 0xFF)      # #FFFFFF
BLACK = (0x00, 0x00, 0x00)      # #000000

# Common theme colors and their approximate RGB values
THEME_COLOR_APPROX = {
    'accent1': (0x44, 0x72, 0xC4),  # Blue
    'accent2': (0xED, 0x7D, 0x31),  # Orange
    'accent3': (0xA5, 0xA5, 0xA5),  # Gray
    'accent4': (0xFF, 0xC0, 0x00),  # Gold
    'accent5': (0x5B, 0x9B, 0xD5),  # Light blue (often fails!)
    'accent6': (0x70, 0xAD, 0x47),  # Green
    'bg1': (0xF2, 0xF2, 0xF2),      # Off-white (often fails!)
    'bg2': (0xE7, 0xE6, 0xE6),      # Light gray
    'tx1': (0x00, 0x00, 0x00),      # Black
    'tx2': (0x44, 0x54, 0x6A),      # Dark gray-blue
    'lt1': (0xFF, 0xFF, 0xFF),      # White
    'dk1': (0x00, 0x00, 0x00),      # Black
}



def relative_luminance(rgb):
    """Calculate relative luminance per WCAG 2.1."""
    def adjust(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * adjust(r) + 0.7152 * adjust(g) + 0.0722 * adjust(b)


def contrast_ratio(color1, color2):
    """Calculate contrast ratio between two RGB colors."""
    l1 = relative_luminance(color1)
    l2 = relative_luminance(color2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def get_text_color(run):
    """Extract RGB color from a text run."""
    try:
        if run.font.color and run.font.color.rgb:
            rgb = run.font.color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except:
        pass

    # Try to get theme color
    try:
        if run.font.color and run.font.color.theme_color:
            theme = str(run.font.color.theme_color).split('(')[0].lower()
            for key, val in THEME_COLOR_APPROX.items():
                if key in theme.lower():
                    return val
    except:
        pass

    # Check XML for scheme color
    try:
        xml = run._r.xml
        match = re.search(r'schemeClr val="([^"]+)"', xml)
        if match:
            scheme = match.group(1).lower()
            if scheme in THEME_COLOR_APPROX:
                return THEME_COLOR_APPROX[scheme]
    except:
        pass

    return None


def get_background_color(shape, slide):
    """Try to determine the background color behind a shape."""
    # Check if shape has its own fill
    try:
        if shape.fill and shape.fill.type and shape.fill.fore_color:
            if shape.fill.fore_color.rgb:
                rgb = shape.fill.fore_color.rgb
                return (rgb[0], rgb[1], rgb[2])
    except:
        pass

    # Look for rectangle shapes that might be backgrounds
    shape_left = shape.left if hasattr(shape, 'left') else 0
    shape_top = shape.top if hasattr(shape, 'top') else 0

    for other in slide.shapes:
        if other == shape:
            continue
        if 'Rectangle' in other.name or 'Background' in other.name:
            try:
                # Check if this rectangle is behind our shape
                other_xml = other._element.xml
                match = re.search(r'srgbClr val="([A-Fa-f0-9]{6})"', other_xml)
                if match:
                    hex_color = match.group(1)
                    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

                # Check for scheme color
                match = re.search(r'schemeClr val="([^"]+)"', other_xml)
                if match:
                    scheme = match.group(1).lower()
                    if scheme in THEME_COLOR_APPROX:
                        return THEME_COLOR_APPROX[scheme]
            except:
                pass

    # Default assumption: white background
    return WHITE


def is_large_text(run):
    """Check if text qualifies as 'large' for WCAG (18pt+ or 14pt+ bold)."""
    try:
        size_pt = run.font.size.pt if run.font.size else 12
        is_bold = run.font.bold
        return size_pt >= 18 or (size_pt >= 14 and is_bold)
    except:
        return False


def analyze_slide(slide, slide_num):
    """Analyze a slide for contrast issues."""
    issues = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text.strip()
                if not text:
                    continue

                text_color = get_text_color(run)
                bg_color = get_background_color(shape, slide)

                if text_color is None:
                    # Can't determine color, skip
                    continue

                ratio = contrast_ratio(text_color, bg_color) if bg_color else None
                required = CONTRAST_LARGE if is_large_text(run) else CONTRAST_NORMAL

                if ratio and ratio < required:
                    issues.append({
                        'slide': slide_num,
                        'shape': shape.name,
                        'text': text[:40] + ('...' if len(text) > 40 else ''),
                        'text_color': text_color,
                        'bg_color': bg_color,
                        'ratio': ratio,
                        'required': required,
                        'is_large': is_large_text(run),
                        'run': run,
                        'shape_obj': shape,
                    })

    return issues


def fix_contrast_issue(issue, slide):
    """Fix a contrast issue by adjusting colors."""
    text_color = issue['text_color']
    bg_color = issue['bg_color']
    run = issue['run']
    shape = issue['shape_obj']

    # Determine if text is light or dark
    text_luminance = relative_luminance(text_color)
    bg_luminance = relative_luminance(bg_color)

    # Strategy: If text is light on light-ish background, darken background
    # If text is dark on dark-ish background, lighten background or darken text more

    if text_luminance > 0.5:  # Light text (like white/off-white)
        # Try to find and darken the background
        bg_fixed = False
        for other_shape in slide.shapes:
            if 'Rectangle' in other_shape.name:
                try:
                    # Check if this might be the background
                    other_shape.fill.solid()
                    other_shape.fill.fore_color.rgb = RGBColor(*DARK_BLUE)
                    bg_fixed = True
                    break
                except:
                    pass

        # Also ensure text is pure white (not off-white)
        run.font.color.rgb = RGBColor(*WHITE)
        fix_type = 'darkened_background' if bg_fixed else 'whitened_text'

    else:  # Dark text
        # Make text darker or change to black
        run.font.color.rgb = RGBColor(*BLACK)
        fix_type = 'darkened_text'

    new_text_color = get_text_color(run) or text_color
    new_bg_color = get_background_color(shape, slide)
    new_ratio = contrast_ratio(new_text_color, new_bg_color) if new_bg_color else None
    required_ratio = issue.get("required", CONTRAST_NORMAL)
    guaranteed = bool(new_ratio and new_ratio >= required_ratio)

    # If still not compliant, force a high-contrast text box fill.
    if new_ratio is None or new_ratio < required_ratio:
        # Prefer black text on white fill.
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*WHITE)
            run.font.color.rgb = RGBColor(*BLACK)
            fix_type = 'forced_white_bg_black_text'
            new_text_color = BLACK
            new_bg_color = WHITE
            new_ratio = contrast_ratio(new_text_color, new_bg_color)
            guaranteed = new_ratio >= required_ratio
        except Exception:
            guaranteed = False

    if not guaranteed and (new_ratio is None or new_ratio < required_ratio):
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*BLACK)
            run.font.color.rgb = RGBColor(*WHITE)
            fix_type = 'forced_black_bg_white_text'
            new_text_color = WHITE
            new_bg_color = BLACK
            new_ratio = contrast_ratio(new_text_color, new_bg_color)
        except Exception:
            pass

    return fix_type


def process_pptx(input_path, output_path=None, analyze_only=False):
    """Process a PPTX file for contrast issues."""
    input_path = Path(input_path)

    if output_path is None:
        output_path = input_path.parent / f"{input_path.stem}_accessible{input_path.suffix}"
    else:
        output_path = Path(output_path)

    prs = Presentation(input_path)

    all_issues = []
    for i, slide in enumerate(prs.slides, 1):
        issues = analyze_slide(slide, i)
        all_issues.extend(issues)

    if not all_issues:
        print(f"  No contrast issues found!")
        return output_path, 0

    print(f"  Found {len(all_issues)} potential contrast issue(s):")
    for issue in all_issues:
        text_hex = '#{:02X}{:02X}{:02X}'.format(*issue['text_color'])
        bg_hex = '#{:02X}{:02X}{:02X}'.format(*issue['bg_color']) if issue['bg_color'] else 'unknown'
        status = "✗" if issue['ratio'] < issue['required'] else "✓"
        print(f"    Slide {issue['slide']}, {issue['shape']}: '{issue['text']}'")
        print(f"      {status} Contrast: {issue['ratio']:.2f}:1 (need {issue['required']}:1)")
        print(f"      Text: {text_hex}, Background: {bg_hex}")

    if analyze_only:
        return None, len(all_issues)

    # Fix issues
    fixes = 0
    for issue in all_issues:
        slide = prs.slides[issue['slide'] - 1]
        fix_type = fix_contrast_issue(issue, slide)
        fixes += 1
        print(f"    Fixed: {fix_type}")

    prs.save(output_path)
    return output_path, fixes


def main():
    parser = argparse.ArgumentParser(
        description='Fix PPTX contrast issues for WCAG 2.1 AA accessibility',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python fix_pptx_contrast.py presentation.pptx
  python fix_pptx_contrast.py presentation.pptx --analyze
  python fix_pptx_contrast.py *.pptx
        """
    )
    parser.add_argument('files', nargs='+', help='PPTX file(s) to process')
    parser.add_argument('-o', '--output', help='Output file (single file only)')
    parser.add_argument('-a', '--analyze', action='store_true',
                        help='Analyze only, do not fix')

    args = parser.parse_args()

    if args.output and len(args.files) > 1:
        print("Error: -o/--output can only be used with a single input file")
        return 1

    print("=" * 60)
    print("PPTX Contrast Fixer")
    print("=" * 60)

    total_fixes = 0

    for file_path in args.files:
        file_path = Path(file_path)

        if not file_path.exists():
            print(f"\nSkipping {file_path} (not found)")
            continue

        if not file_path.suffix.lower() == '.pptx':
            print(f"\nSkipping {file_path} (not a .pptx file)")
            continue

        print(f"\nProcessing: {file_path.name}")

        try:
            output_path, num_fixes = process_pptx(
                file_path,
                args.output,
                analyze_only=args.analyze
            )
            total_fixes += num_fixes

            if output_path and not args.analyze:
                print(f"  ✓ Saved: {output_path.name}")

        except Exception as e:
            print(f"  ✗ Error: {e}")

    print("\n" + "=" * 60)
    if args.analyze:
        print(f"Analysis complete. Found {total_fixes} potential issue(s).")
    else:
        print(f"Complete! Fixed {total_fixes} issue(s).")
    print("=" * 60)

    return 0


if __name__ == "__main__":
    exit(main())
